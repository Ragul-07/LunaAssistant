from flask import Flask, request, jsonify
from flask_cors import CORS
import pyttsx3
import os
import webbrowser
from docx import Document
from pptx import Presentation
import pyautogui
import datetime
import pandas as pd
import random
import time  # Import time module to introduce delays

app = Flask(__name__)
CORS(app)

# Initialize voice preference
current_voice = "Female"

@app.route('/api/set-voice', methods=['POST'])
def set_voice():
    global current_voice
    data = request.get_json()
    if 'voice' in data:
        current_voice = data['voice']
        return jsonify({"message": "Voice updated successfully", "current_voice": current_voice}), 200
    return jsonify({"error": "Invalid request"}), 400

def speak(text):
    global current_voice
    engine = pyttsx3.init('sapi5')  # Initialize engine inside the function
    voices = engine.getProperty('voices')

    # Set voice based on current preference
    if current_voice == "Male":
        engine.setProperty('voice', voices[0].id)  # Male voice
    else:
        engine.setProperty('voice', voices[1].id)  # Female voice

    engine.setProperty('rate', 170)
    engine.say(text)
    engine.runAndWait()

@app.route('/api/speak', methods=['POST'])
def speak_text():
    data = request.get_json()
    if 'text' in data:
        speak(data['text'])
        return jsonify({"message": "Speaking: " + data['text']}), 200
    return jsonify({"error": "No text provided"}), 400

def open_application(app_name):
    app_paths = {
        "notepad": "notepad.exe",
        "calculator": "calc.exe",
        "chrome": "start cmd /c chrome.exe",
        "word": "winword.exe",
        "excel": "excel.exe",
        "powerpoint": "powerpnt.exe",
        "vscode": "start cmd /c code",
        "paint": "mspaint.exe",
        "explorer": "explorer.exe"
    }
    if app_name in app_paths:
        os.system(f"start {app_paths[app_name]}")
        speak(f"Opening {app_name} successfully")
        return f"Opened {app_name} successfully"
    return "Application not found"

def search_on_google(query):
    url = f"https://www.google.com/search?q={query}"
    webbrowser.open(url)
    speak(f"Searching Google for {query} successfully")
    return f"Searching Google for {query} successfully"

def play_youtube(video):
    url = f"https://www.youtube.com/results?search_query={video}"
    webbrowser.open(url)
    speak(f"Playing {video} on YouTube successfully")
    return f"Playing {video} on YouTube successfully"

def get_date_time():
    now = datetime.datetime.now()
    date_time = now.strftime("%Y-%m-%d %H:%M:%S")
    return f"Current date and time: {date_time}"

def create_ppt(topic):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = topic.capitalize()
    slide.shapes.placeholders[1].text = "Presentation by Luna"
    topics = ["Introduction", "Features", "Future", "Conclusion"]
    for t in topics:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = t
        slide.placeholders[1].text = f"Details about {t}"
    ppt.save(f"{topic}.pptx")
    os.system(f"start {topic}.pptx")
    speak(f"Presentation on {topic} created successfully")
    return f"Created presentation on {topic} successfully"

def create_letter(reason, recipient):
    doc = Document()
    doc.add_heading(f"Letter for {reason}", 0)
    doc.add_paragraph(f"Dear {recipient},\n\nI am writing to discuss {reason}.\n\nBest Regards,\nUser")
    doc.save(f"Letter_{reason}.docx")
    os.system(f"start Letter_{reason}.docx")
    speak(f"Letter for {reason} created successfully")
    return f"Letter for {reason} created successfully"

def create_excel():
    data = {
        "Name": ["Alice", "Bob", "Charlie"],
        "Age": [25, 30, 35],
        "City": ["New York", "Los Angeles", "Chicago"]
    }
    df = pd.DataFrame(data)
    df.to_excel("data.xlsx", index=False)
    os.system("start data.xlsx")
    speak("Excel sheet created successfully")
    return "Excel sheet created successfully"

def text_to_speech(text):
    filename = "speech_output.mp3"
    engine.save_to_file(text, filename)
    engine.runAndWait()
    os.system(f"start {filename}")
    speak("Text converted to speech successfully")
    return "Text converted to speech successfully"

def open_vscode_and_create_java_file():
    os.system("start cmd /c code")  # Launch VS Code using CMD
    speak("Opening VS Code, please wait")
    time.sleep(5)

    pyautogui.hotkey('alt', 'tab')  # Switch to VS Code
    time.sleep(3)
    pyautogui.hotkey('win', 'up')  # Maximize the window
    time.sleep(3)

    pyautogui.hotkey('ctrl', 'n')  
    time.sleep(4)

    java_code = """
    public class Pattern {
        public static void main(String[] args) {
            int n = 5;
            for (int i = 1; i <= n; i++) {
                for (int j = 1; j <= i; j++) {
                    System.out.print("* ");
                }
                System.out.println();
            }
        }
    }
    """

    pyautogui.write(java_code, interval=0.05)
    time.sleep(1)

    pyautogui.hotkey('ctrl', 's')  
    time.sleep(1)
    pyautogui.write("Pattern.java")  
    pyautogui.press('enter')
    time.sleep(1)

    speak("Java pattern program file created successfully")
    return "Java pattern program file created successfully"
def write_letter_from_command(command):
    # Extract reason and recipient from the command
    if "write a" in command and "letter to" in command:
        parts = command.split("letter to")
        reason_part = parts[0].replace("open word and write a", "").strip()
        recipient = parts[1].strip()
        reason = reason_part if reason_part else "General"

        # Create the letter document
        return create_letter(reason, recipient)
    
    return "Invalid command format"




def get_time_based_greeting():
    """ Returns a greeting based on the current time of day """
    current_hour = datetime.datetime.now().hour
    if current_hour < 12:
        return "Good morning! How can I assist you today?"
    elif current_hour < 18:
        return "Good afternoon! What do you need help with?"
    else:
        return "Good evening! How can I make your night easier?"

def chatbot_response(user_input):
    """ AI Chatbot function with voice response """
    user_input = user_input.lower()

    responses = {
        ("hello", "hi", "hey", "greetings", "what's up"): get_time_based_greeting(),
        ("how are you", "how do you feel", "are you okay"): "I'm just a virtual assistant, but I'm always here for you! How about you?",
        ("what is your name", "who are you", "tell me your name"): "I'm Luna, your AI assistant! What's your name?",
        ("who created you", "who made you", "who is your developer"): "I was created by an awesome developer who wants to make life easier for you!",
        
        # Motivational
        ("i'm feeling sad", "cheer me up", "i feel down"): random.choice([
            "I'm here for you. Want to talk about it?",
            "You're stronger than you think! Want to hear a joke?",
            "Bad days happen, but they don’t define you! You got this!"
        ]),
        ("tell me something motivational", "motivate me", "inspire me"): random.choice([
            "Success is not final, failure is not fatal: It is the courage to continue that counts.",
            "Every day is a fresh start. Keep going!",
            "You are capable of amazing things. Keep pushing forward!"
        ]),

        # Fun & Jokes
        ("tell me a joke", "make me laugh", "say something funny"): random.choice([
            "Why don’t skeletons fight each other? Because they don’t have the guts!",
            "Parallel lines have so much in common. It’s a shame they’ll never meet.",
            "Why did the computer catch a cold? Because it left its Windows open!"
        ]),

        # Facts
        ("tell me a fun fact", "random fact", "something interesting"): random.choice([
            "Did you know honey never spoils? Archaeologists found pots over 3000 years old and still edible!",
            "Octopuses have three hearts, and their blood is blue!",
            "Bananas are berries, but strawberries aren't!"
        ]),

        # Technology
        ("what is artificial intelligence", "explain ai", "ai definition"): "Artificial Intelligence (AI) is the simulation of human intelligence in machines that can learn, reason, and make decisions.",
        ("who owns apple", "who founded apple", "apple company history"): "Apple was founded by Steve Jobs, Steve Wozniak, and Ronald Wayne in 1976.",

        # Productivity
        ("how do i be productive", "increase productivity", "stay focused"): "Start with small tasks, avoid multitasking, take breaks, and set clear goals. You've got this!",
    }

    # **Matching Logic**
    for keys, response in responses.items():
        if any(keyword in user_input for keyword in keys):
            speak(response)  # Speak the response
            return response

    default_response = "I'm here to chat! Ask me anything!"
    speak(default_response)  # Speak default response
    return default_response


@app.route("/command", methods=["POST"])
def handle_command():
    data = request.get_json()
    command = data.get("command", "").lower().strip()

    if "open" in command and "write a pattern program in java" in command:
        open_application("vscode")
        response = open_vscode_and_create_java_file()
    elif command.startswith("open "):
        app_name = command[5:].strip()
        response = open_application(app_name)
    elif any(keyword in command for keyword in ["search", "google", "surf the internet"]):
        query = command.split("search", 1)[-1].strip()
        response = search_on_google(query)
    elif any(keyword in command for keyword in ["play", "play on youtube", "watch youtube"]):
        video = command.split("play", 1)[-1].strip()
        response = play_youtube(video)
    elif any(keyword in command for keyword in ["date", "time"]):
        response = get_date_time()
    elif any(keyword in command for keyword in ["make a ppt", "create presentation", "generate slides"]):
        topic = command.split("on")[-1].strip() if "on" in command else "Untitled"
        response = create_ppt(topic)
    elif any(keyword in command for keyword in ["write a letter", "make a document", "open word and assist me"]):
        reason = command.split("letter")[-1].strip() if "letter" in command else "General"
        response = create_letter(reason, "Recipient")
    elif any(keyword in command for keyword in ["create excel", "generate a sample excel file"]):
        response = create_excel()
    elif "convert to speech" in command:
        text = command.split("convert to speech", 1)[-1].strip()
        response = text_to_speech(text)
    else:
        response = chatbot_response(command)

    return jsonify({"response": response})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)      